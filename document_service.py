from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Iterable

import pdf_service as pdfs


@dataclass
class DocumentChunk:
    label: str
    text: str
    score: int = 0


WORK_TERMS = [
    *pdfs.HERWANTO_TERMS,
    *pdfs.TIMETABLE_TERMS,
    "nbss",
    "naval base",
    "school",
    "deadline",
    "submission",
    "meeting",
    "briefing",
    "event",
    "lesson",
    "worksheet",
    "assessment",
    "rubric",
    "proposal",
    "cca",
    "football",
    "gameplan",
    "ruh",
]


def _clean_text(value: str) -> str:
    lines = [line.strip() for line in (value or "").splitlines()]
    return "\n".join(line for line in lines if line)


def _score_text(text: str, extra_terms: Iterable[str] = ()) -> int:
    lower = (text or "").lower()
    terms = [t.lower() for t in [*WORK_TERMS, *extra_terms] if str(t).strip()]
    score = 0
    for term in terms:
        if term in lower:
            score += 8 if term in pdfs.HERWANTO_TERMS else 2
    score += min(lower.count("\n"), 120) // 10
    score += len(re.findall(r"\b\d{1,2}[:.]\d{2}\b|\b\d{1,2}/\d{1,2}/\d{2,4}\b|\b202\d\b", lower))
    return score


def _caption_terms(caption: str) -> list[str]:
    return re.findall(r"[a-z0-9_.@-]+", (caption or "").lower())


def extract_pdf(file_bytes: bytes, caption: str = "") -> tuple[str, str]:
    page_count, pages = pdfs.extract_pdf_pages(file_bytes)
    excerpt, selected_pages, text_pages = pdfs.build_pdf_excerpt(pages, caption=caption)
    index = pdfs.format_pdf_index(page_count, text_pages, selected_pages)
    return index, excerpt


def extract_docx(file_bytes: bytes, caption: str = "") -> tuple[str, str]:
    from docx import Document

    document = Document(io.BytesIO(file_bytes))
    chunks: list[DocumentChunk] = []
    paragraphs = [_clean_text(p.text) for p in document.paragraphs if _clean_text(p.text)]
    if paragraphs:
        chunks.append(DocumentChunk("Paragraphs", "\n".join(paragraphs)))

    for t_idx, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [_clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(DocumentChunk(f"Table {t_idx}", "\n".join(rows)))

    terms = _caption_terms(caption)
    for chunk in chunks:
        chunk.score = _score_text(chunk.text, terms)
    ranked = sorted(chunks, key=lambda c: (c.score, len(c.text)), reverse=True)
    selected = sorted(ranked[:16], key=lambda c: chunks.index(c))
    excerpt = "\n\n".join(
        f"--- {chunk.label} (relevance {chunk.score}) ---\n{chunk.text[:4500]}"
        for chunk in selected
    )
    index = f"DOCX has {len(paragraphs)} paragraphs and {len(document.tables)} tables; analysed {len(selected)} chunks."
    return index, excerpt


def extract_pptx(file_bytes: bytes, caption: str = "") -> tuple[str, str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: list[DocumentChunk] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and _clean_text(shape.text):
                texts.append(_clean_text(shape.text))
        notes = ""
        try:
            notes = _clean_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes = ""
        if notes:
            texts.append(f"Notes:\n{notes}")
        if texts:
            chunks.append(DocumentChunk(f"Slide {idx}", "\n".join(texts)))

    terms = _caption_terms(caption)
    for chunk in chunks:
        chunk.score = _score_text(chunk.text, terms)
    ranked = sorted(chunks, key=lambda c: (c.score, len(c.text)), reverse=True)
    selected = sorted(ranked[:20], key=lambda c: chunks.index(c))
    excerpt = "\n\n".join(
        f"--- {chunk.label} (relevance {chunk.score}) ---\n{chunk.text[:3500]}"
        for chunk in selected
    )
    index = f"PPTX has {len(prs.slides)} slides; analysed {len(selected)} slides."
    return index, excerpt


def extract_supported_document(file_bytes: bytes, mime_type: str, filename: str = "", caption: str = "") -> tuple[str, str, str]:
    mime = (mime_type or "").lower()
    name = (filename or "").lower()
    if mime == "application/pdf" or name.endswith(".pdf"):
        index, excerpt = extract_pdf(file_bytes, caption)
        return "PDF", index, excerpt
    if (
        mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or name.endswith(".docx")
    ):
        index, excerpt = extract_docx(file_bytes, caption)
        return "DOCX", index, excerpt
    if (
        mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or name.endswith(".pptx")
    ):
        index, excerpt = extract_pptx(file_bytes, caption)
        return "PPTX", index, excerpt
    raise ValueError(f"Unsupported document type: {mime_type or filename}")
