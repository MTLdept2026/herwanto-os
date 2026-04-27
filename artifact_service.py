from __future__ import annotations

import re
from pathlib import Path
from typing import Any


ARTIFACT_DIR = Path("files/generated")


def _safe_filename(value: str, fallback: str) -> str:
    clean = re.sub(r"[^A-Za-z0-9._-]+", "-", (value or "").strip()).strip("-._")
    return clean[:80] or fallback


def _clean_text(value: Any) -> str:
    return str(value or "").strip()


def _normalise_bullets(value: Any) -> list[str]:
    if isinstance(value, list):
        return [_clean_text(item) for item in value if _clean_text(item)]
    text = _clean_text(value)
    return [text] if text else []


def render_docx(spec: dict, prefix: str = "document") -> Path:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    title = _clean_text(spec.get("title")) or "Untitled Document"
    path = ARTIFACT_DIR / f"{_safe_filename(title, prefix)}.docx"

    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(11)

    title_para = document.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = _clean_text(spec.get("subtitle"))
    if subtitle:
        para = document.add_paragraph(subtitle)
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    author = _clean_text(spec.get("author"))
    if author:
        para = document.add_paragraph(author)
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for section in spec.get("sections", []):
        if not isinstance(section, dict):
            continue
        heading = _clean_text(section.get("heading"))
        if heading:
            document.add_heading(heading, level=1)

        body = _clean_text(section.get("body"))
        if body:
            for paragraph in [p.strip() for p in body.split("\n") if p.strip()]:
                document.add_paragraph(paragraph)

        for bullet in _normalise_bullets(section.get("bullets")):
            document.add_paragraph(bullet, style="List Bullet")

    document.save(path)
    return path


def render_pptx(spec: dict, prefix: str = "slides") -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    title = _clean_text(spec.get("title")) or "Untitled Deck"
    path = ARTIFACT_DIR / f"{_safe_filename(title, prefix)}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    subtitle = title_slide.placeholders[1]
    subtitle.text = _clean_text(spec.get("subtitle")) or _clean_text(spec.get("audience"))

    for slide_spec in spec.get("slides", []):
        if not isinstance(slide_spec, dict):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = _clean_text(slide_spec.get("title")) or "Slide"
        body = slide.placeholders[1].text_frame
        body.clear()

        bullets = _normalise_bullets(slide_spec.get("bullets"))
        if not bullets and _clean_text(slide_spec.get("body")):
            bullets = [_clean_text(slide_spec.get("body"))]

        for idx, bullet in enumerate(bullets[:7]):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(24 if len(bullets) <= 4 else 20)

        notes = _clean_text(slide_spec.get("notes"))
        if notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = notes

    prs.save(path)
    return path
